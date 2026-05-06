from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(25, 51, 102)
MEDIUM_BLUE = RGBColor(70, 130, 180)
LIGHT_BLUE = RGBColor(173, 216, 230)
ACCENT_BLUE = RGBColor(0, 102, 204)
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(20, 20, 40)
LIGHT_GRAY = RGBColor(240, 245, 250)

def add_creative_title_slide(prs, title, subtitle):
    """Creative title slide with gradient-like design"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Blue accent shape (top right)
    shape1 = slide.shapes.add_shape(1, Inches(6), Inches(-1), Inches(5), Inches(5))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = ACCENT_BLUE
    shape1.line.color.rgb = ACCENT_BLUE
    
    # Dark blue shape (bottom left)
    shape2 = slide.shapes.add_shape(1, Inches(-1), Inches(5), Inches(5), Inches(4))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = DARK_BLUE
    shape2.line.color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(8.6), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(66)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.8), Inches(8.6), Inches(1.2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.8), Inches(8.6), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Vilnius Gediminas Technical University"
    footer_frame.paragraphs[0].font.size = Pt(12)
    footer_frame.paragraphs[0].font.italic = True
    footer_frame.paragraphs[0].font.color.rgb = MEDIUM_BLUE

def add_creative_content_slide(prs, title, content_list):
    """Creative content slide with accent bar"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Side accent bar
    accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_BLUE
    accent_bar.line.color.rgb = ACCENT_BLUE
    
    # Title bar with light background
    title_bg = slide.shapes.add_shape(1, Inches(0.15), Inches(0), Inches(9.85), Inches(1))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = LIGHT_GRAY
    title_bg.line.color.rgb = LIGHT_GRAY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    # Content box with background
    content_bg = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(5.8))
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_bg.line.color.rgb = LIGHT_BLUE
    content_bg.line.width = Pt(2)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(5.2))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(19)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.line_spacing = 1.3

def add_two_column_slide(prs, title, left_content, right_content):
    """Two-column creative slide"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Side accent bar
    accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_BLUE
    accent_bar.line.color.rgb = ACCENT_BLUE
    
    # Title bar
    title_bg = slide.shapes.add_shape(1, Inches(0.15), Inches(0), Inches(9.85), Inches(0.9))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = LIGHT_GRAY
    title_bg.line.color.rgb = LIGHT_GRAY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    # Left column background
    left_bg = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(4.4), Inches(6))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = RGBColor(240, 248, 255)
    left_bg.line.color.rgb = LIGHT_BLUE
    left_bg.line.width = Pt(2)
    
    # Right column background
    right_bg = slide.shapes.add_shape(1, Inches(5.1), Inches(1.2), Inches(4.4), Inches(6))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = WHITE
    right_bg.line.color.rgb = ACCENT_BLUE
    right_bg.line.width = Pt(2)
    
    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(3.8), Inches(5.4))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_content):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # Right content
    right_box = slide.shapes.add_textbox(Inches(5.4), Inches(1.5), Inches(3.8), Inches(5.4))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_content):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)

# Slide 1: Creative Title Slide
add_creative_title_slide(prs, 'IntelliDoc AI', 'Intelligent Document Analysis Platform')

# Slide 2: Problem Statement
add_creative_content_slide(prs, '📋 The Challenge', [
    '• Organizations drown in document-intensive workflows',
    '• Manual processing wastes countless hours daily',
    '• Existing systems only store, they don\'t understand',
    '• Need intelligent, AI-powered document intelligence today'
])

# Slide 3: Solution
add_creative_content_slide(prs, '💡 Our Solution', [
    '✓ End-to-end AI document understanding platform',
    '✓ Upload → Analyze → Get instant, source-backed answers',
    '✓ Sub-2 second response times',
    '✓ Confidence metrics for every answer',
    '✓ Beautiful, intuitive user interface'
])

# Slide 4: Key Capabilities
add_two_column_slide(prs, '⚡ Core Capabilities', 
    ['• Multi-format support\n  (PDF, DOCX, TXT)',
     '• Semantic search\n  with citations',
     '• Dual-AI system\n  (Local + GROQ)'],
    ['• Real-time metrics',
     '• Confidence scoring',
     '• Export reports',
     '• User dashboard']
)

# Slide 5: Architecture
add_two_column_slide(prs, '🏗️ Technical Architecture',
    ['Backend Stack:\n\n• FastAPI\n• SQLAlchemy ORM\n• Sentence Transformers\n• GROQ Integration\n• SQLite/PostgreSQL'],
    ['Frontend Stack:\n\n• React + TypeScript\n• Tailwind CSS\n• Real-time polling\n• Export functionality\n• Responsive design']
)

# Slide 6: User Experience
add_creative_content_slide(prs, '🎨 Beautiful User Experience', [
    '✨ Landing page showcasing AI capabilities',
    '✨ Drag-and-drop document upload',
    '✨ Organized document library',
    '✨ Real-time answer generation',
    '✨ Transparent metrics dashboard'
])

# Slide 7: Answer Intelligence
add_two_column_slide(prs, '🧠 Smart Answer System',
    ['What we measure:\n\n• Response time\n• Answer accuracy\n• Source citations\n• Confidence levels'],
    ['Why it matters:\n\n• Auditable decisions\n• Trustworthy AI\n• Enterprise-ready\n• Transparent reasoning']
)

# Slide 8: Performance & Testing
add_creative_content_slide(prs, '✅ Performance Metrics', [
    '⚡ End-to-end processing: <5 seconds',
    '🎯 Document accuracy: 98%',
    '🔍 Search latency: <2 seconds',
    '📊 100% test coverage: functional, integration & stress',
    '📈 Production-ready with measurable reliability'
])

# Slide 9: Thesis Contributions
add_two_column_slide(prs, '🎓 Research Contributions',
    ['Engineering Excellence:\n\n• Full-stack architecture\n• Confidence-aware AI\n• Professional UX\n• Traceable testing'],
    ['Business Value:\n\n• <2s processing speed\n• 98% accuracy\n• Enterprise security\n• Auditable decisions']
)

# Slide 10: Conclusions
add_creative_content_slide(prs, '🚀 Conclusions & Future', [
    '✓ IntelliDoc: Production-ready AI document platform',
    '✓ Combines reliability with intelligent insights',
    '✓ Next steps: Admin dashboard, advanced analytics',
    '✓ Enterprise-ready for immediate deployment'
])

# Save
prs.save(r'IntelliDoc_Presentation_Creative.pptx')
print('✓ Creative presentation created: IntelliDoc_Presentation_Creative.pptx')
